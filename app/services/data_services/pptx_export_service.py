import io
import re
from typing import Dict, List

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _clean_text(text: str) -> str:
    text = re.sub(r"```[\s\S]*?```", "代码/图示内容见原资源正文。", text or "")
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"#+\s*", "", text)
    return text.strip()


def _split_sections(content: str) -> List[Dict[str, List[str]]]:
    lines = (content or "").splitlines()
    sections = []
    current = {"title": "核心内容", "items": []}

    for line in lines:
        raw = line.strip()
        if not raw:
            continue

        heading = re.match(r"^#{1,3}\s+(.+)$", raw)
        if heading:
            if current["items"]:
                sections.append(current)
            current = {"title": _clean_text(heading.group(1))[:40], "items": []}
            continue

        cleaned = _clean_text(re.sub(r"^[-*]\s+|^\d+[.、)]\s+", "", raw))
        if cleaned:
            current["items"].append(cleaned[:120])

    if current["items"]:
        sections.append(current)

    if not sections:
        plain = _clean_text(content)
        chunks = [plain[i:i + 90] for i in range(0, len(plain), 90)]
        sections.append({"title": "核心内容", "items": chunks[:6] or ["暂无正文内容"]})

    return sections[:8]


def _set_run_font(run, size=22, bold=False):
    run.font.name = "PingFang SC"
    run.font.size = Pt(size)
    run.font.bold = bold


def _add_title(slide, title: str, subtitle: str = ""):
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(8.0), Inches(0.8))
    frame = title_box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_run_font(run, 30, True)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.0), Inches(0.45))
        sub_frame = sub_box.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        _set_run_font(run, 14, False)


def _add_bullets(slide, items: List[str]):
    body = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(8.1), Inches(4.55))
    frame = body.text_frame
    frame.word_wrap = True
    frame.clear()

    for index, item in enumerate(items[:7]):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.level = 0
        p.text = f"• {item}"
        p.font.name = "PingFang SC"
        p.font.size = Pt(18)


def build_resource_pptx(resource: Dict):
    title = resource.get("title") or "LingXi 学习资源"
    resource_type = resource.get("type") or "学习资源"
    summary = resource.get("summary") or "由灵析多智能体学习系统生成。"
    content = resource.get("content") or ""
    source = resource.get("source") or resource.get("uploader") or "LingXi 知识库"

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    cover = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(cover, title[:34], f"{resource_type} · {source}")
    summary_box = cover.shapes.add_textbox(Inches(0.75), Inches(2.05), Inches(8.3), Inches(1.4))
    summary_frame = summary_box.text_frame
    summary_frame.word_wrap = True
    summary_frame.clear()
    p = summary_frame.paragraphs[0]
    run = p.add_run()
    run.text = summary[:180]
    _set_run_font(run, 20, False)

    outline = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(outline, "学习资源结构", "系统自动提取资源正文结构，生成可直接演示的课件。")
    sections = _split_sections(content)
    _add_bullets(outline, [section["title"] for section in sections])

    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title(slide, section["title"])
        _add_bullets(slide, section["items"])

    final = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(final, "学习建议", "结合画像、路径规划与资源审核结果继续学习。")
    _add_bullets(final, [
        "先阅读文字讲解，建立概念框架。",
        "再查看流程图或代码注释，理解知识之间的关系。",
        "完成分步题解和实践任务，记录错因与改进点。",
        "如资源来自 AI 生成，请结合管理员审核建议进行复核。",
    ])

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    filename = re.sub(r"[^\w\u4e00-\u9fff-]+", "_", title).strip("_") or "lingxi_resource"
    return output, f"{filename}.pptx"
